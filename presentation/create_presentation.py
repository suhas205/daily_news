from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re

def create_title_slide(prs, title, subtitle=""):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    return slide

def create_content_slide(prs, title, content):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = title
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for line in content:
        p = tf.add_paragraph()
        p.text = line.strip('- ')
        p.level = line.count('  ')  # Count indentation
    
    return slide

def create_code_slide(prs, title, code):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Add code
    code_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    code_box.text = code
    code_box.text_frame.paragraphs[0].font.name = 'Courier New'
    code_box.text_frame.paragraphs[0].font.size = Pt(14)
    
    return slide

def parse_markdown_and_create_pptx(md_file, pptx_file):
    prs = Presentation()
    
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split content into slides
    slides = content.split('---')
    
    for slide_content in slides:
        if not slide_content.strip():
            continue
            
        lines = [line for line in slide_content.split('\n') if line.strip()]
        
        if not lines:
            continue
            
        # Get title (remove # symbols)
        title = lines[0].lstrip('#').strip()
        content_lines = lines[1:]
        
        # Check if this is a code slide
        if '```' in slide_content:
            code_start = slide_content.find('```') + 3
            code_end = slide_content.find('```', code_start)
            if code_end != -1:
                code = slide_content[code_start:code_end].strip()
                create_code_slide(prs, title, code)
        else:
            # If it's the first slide, create a title slide
            if slides.index(slide_content) == 0:
                subtitle = content_lines[0].lstrip('#').strip() if content_lines else ""
                create_title_slide(prs, title, subtitle)
            else:
                create_content_slide(prs, title, content_lines)
    
    prs.save(pptx_file)

if __name__ == "__main__":
    md_file = "Sports_News_App_Presentation.md"
    pptx_file = "Sports_News_App_Presentation.pptx"
    parse_markdown_and_create_pptx(md_file, pptx_file)
